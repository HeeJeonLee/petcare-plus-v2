#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PetCare+ 사업계획서 PPT 생성 스크립트
최신 정보: 2026-02-28 (수인AI브릿지, 사업자등록번호: 151-09-03201)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_business_plan_ppt():
    """PetCare+ 사업계획서 PPT 생성"""

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    BLUE = RGBColor(41, 128, 185)
    DARK_BLUE = RGBColor(25, 77, 138)
    LIGHT_BLUE = RGBColor(174, 214, 241)
    ORANGE = RGBColor(230, 126, 34)
    GREEN = RGBColor(46, 204, 113)
    GRAY = RGBColor(52, 73, 94)
    WHITE = RGBColor(255, 255, 255)

    def add_title_slide(title, subtitle):
        """제목 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

        # 배경색
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 부제목
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title, content_list):
        """내용 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 제목 배경
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = BLUE
        title_shape.line.color.rgb = BLUE

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 내용
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_list):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0

    def add_two_column_slide(title, left_title, left_items, right_title, right_items):
        """2열 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 제목 배경
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = BLUE
        title_shape.line.color.rgb = BLUE

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 왼쪽 컬럼
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5.8))
        text_frame = left_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE

        for item in left_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_before = Pt(6)
            p.level = 0

        # 오른쪽 컬럼
        right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(5.8))
        text_frame = right_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE

        for item in right_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_before = Pt(6)
            p.level = 0

    # ========== 슬라이드 생성 ==========

    # 1. 표지
    add_title_slide(
        "🐾 PetCare+",
        "대한민국 최고의 AI 기반 펫보험 비교 플랫폼\n사업계획서\n\n2026년 02월"
    )

    # 2. 회사 개요
    add_content_slide(
        "📋 회사 개요",
        [
            "🏢 회사명: 수인AI브릿지",
            "",
            "📍 주소: 경기도 수원시 영통구",
            "   (상세: 동탄원천로1109번길 37, 103층 502호)",
            "",
            "📊 사업자등록번호: 151-09-03201",
            "",
            "📅 개업일: 2026년 02월 25일",
            "",
            "👤 대표자: 이희천",
            "",
            "📱 연락처: 010-5650-0670",
            "   이메일: hejunl@hanmail.net"
        ]
    )

    # 3. 사업 개요
    add_content_slide(
        "🎯 사업 개요",
        [
            "💡 비전",
            "  대한민국 최고의 AI 기반 펫보험 비교 플랫폼",
            "",
            "🎁 핵심 가치",
            "  • 24시간 무료 AI 상담 (Claude AI 기반)",
            "  • 8개 보험사 실시간 비교",
            "  • 펫보험 가입 전 완벽한 정보 제공",
            "",
            "🚀 사업 모델",
            "  • B2C: 반려동물 보호자 대상 무료 서비스",
            "  • B2B: 보험사 제휴 및 마케팅 파트너십",
            "  • 향후: 광고, 제휴 수수료 기반 수익화"
        ]
    )

    # 4. 핵심 기능 (Part 1)
    add_content_slide(
        "⭐ 핵심 기능 (1/2)",
        [
            "1️⃣ AI 맞춤형 펫보험 추천",
            "  Claude AI를 활용한 개인화된 추천 알고리즘",
            "",
            "2️⃣ 8개 보험사 상세 비교표",
            "  • 메리츠, 삼성, 현대, KB, 한화, DB, 농협, 롯데",
            "  • 보험료, 보장내용, 특징, 리뷰 포함",
            "",
            "3️⃣ AI 챗봇 상담 (24시간)",
            "  Claude API 기반 고급 상담 시스템",
            "  • 펫 건강 정보",
            "  • 보험 선택 가이드"
        ]
    )

    # 5. 핵심 기능 (Part 2)
    add_content_slide(
        "⭐ 핵심 기능 (2/2)",
        [
            "4️⃣ 주변 동물병원 검색",
            "  Google Maps API를 활용한 실시간 병원 검색",
            "  • 위치 기반 검색",
            "  • 거리 표시 및 길찾기",
            "",
            "5️⃣ 보험금 청구 프로세스 가이드",
            "  • 단계별 청구 방법",
            "  • 필요 서류 안내",
            "",
            "6️⃣ 온라인 상담 신청",
            "  실시간 상담 신청 → 전문가 연결"
        ]
    )

    # 6. 시장 분석
    add_two_column_slide(
        "📊 시장 분석",
        "시장 규모",
        [
            "• 국내 펫보험 시장",
            "  2023년: 약 2,500억원",
            "  2026년 예상: 4,000억원+",
            "",
            "• 반려동물 소유자",
            "  약 1,180만 명 (42%)",
            "",
            "• 펫보험 가입률",
            "  약 15-20% (성장중)"
        ],
        "경쟁 우위",
        [
            "✅ AI 기반 맞춤형",
            "   대부분 수동 비교",
            "",
            "✅ 24시간 무료 상담",
            "   오픈 시간 제한 경쟁사",
            "",
            "✅ 8개 보험사 통합",
            "   단일 또는 소수만 지원",
            "",
            "✅ 사용자 친화적 UX"
        ]
    )

    # 7. 기술 스택
    add_two_column_slide(
        "🛠️ 기술 스택",
        "Frontend",
        [
            "• React 18",
            "• Vite (번들러)",
            "• Tailwind CSS",
            "• Google Maps API",
            "• PWA 지원"
        ],
        "Backend & Infrastructure",
        [
            "• Node.js + Express",
            "• Supabase (DB)",
            "• Claude API",
            "• Resend (이메일)",
            "• Vercel (호스팅)"
        ]
    )

    # 8. 수익 모델
    add_content_slide(
        "💰 수익 모델",
        [
            "Phase 1 (현재): 무료 서비스 통한 사용자 확보",
            "",
            "Phase 2 (6개월): 보험사 제휴 수익",
            "  • 제휴 수수료 (가입당 5~10%)",
            "  • 광고 수익 (검색 상단 노출)",
            "  • 보험사 데이터 API 라이센싱",
            "",
            "Phase 3 (1년): 프리미엄 기능",
            "  • 상담 프리미엄 (월 9,900원)",
            "  • 맞춤형 리포트 (월 4,900원)",
            "",
            "예상 연간 매출 (2027년): 5억원~10억원"
        ]
    )

    # 9. 팀 구성
    add_two_column_slide(
        "👥 팀 구성",
        "현재 (MVP)",
        [
            "👤 이희천",
            "  역할: 대표이사",
            "  경력: 보험컨설턴트",
            "          AI 비즈니스 개발",
            "",
            "🤖 Claude AI",
            "  역할: AI 상담사",
            "  역량: 24시간 응답"
        ],
        "향후 확대 계획",
        [
            "👨‍💼 영업팀",
            "  보험사 제휴 담당",
            "",
            "👨‍💻 개발팀",
            "  프로덕션 운영 및 개선",
            "",
            "🎨 UX/UI팀",
            "  사용자 경험 최적화"
        ]
    )

    # 10. 자율진화형 플랫폼
    add_content_slide(
        "🚀 자율진화형 플랫폼 (향후 계획)",
        [
            "📊 자동 데이터 갱신",
            "  • 보험료 일일 자동 수집 및 동기화",
            "  • 변경사항 자동 감지 및 반영",
            "",
            "📝 콘텐츠 자동 생성",
            "  • Claude API로 일일 블로그 포스트 생성",
            "  • SNS 콘텐츠 자동화 (Instagram, Twitter, Facebook)",
            "",
            "🤖 AI 기반 자동 최적화",
            "  • 사용자 행동 분석 (Analytics)",
            "  • 개선안 자동 도출 및 배포",
            "  • A/B 테스트 자동화",
            "",
            "🎯 목표: 90% 자동화 (인간 개입 최소화)"
        ]
    )

    # 11. 시간표 (Roadmap)
    add_content_slide(
        "📅 개발 시간표 (Roadmap)",
        [
            "✅ Phase 1: MVP 완성 (2026.02.28 완료)",
            "  • 6가지 이슈 해결 완료",
            "  • 프로덕션 배포 준비",
            "",
            "🟡 Phase 2: 자율진화형 플랫폼 (2026.03~05)",
            "  • 데이터 자동 갱신 시스템",
            "  • 콘텐츠 자동 생성",
            "  • Analytics 고도화",
            "",
            "🟢 Phase 3: 사업 확장 (2026.06~12)",
            "  • 보험사 제휴 확대",
            "  • 광고 수익 활성화",
            "  • 국제 확장 (영문 지원)"
        ]
    )

    # 12. 법적 준칙
    add_content_slide(
        "⚖️ 법적 준칙 및 준수사항",
        [
            "✅ 면책 공고 3곳 배치",
            "  '본 페이지는 정보 제공이며 상품 권유가 아님'",
            "",
            "✅ 보험업법 준수",
            "  • 확정적 표현 제거 ('최저가', '무조건 보상')",
            "  • 약관 확인 필수 안내",
            "",
            "✅ 개인정보보호법 준수",
            "  • Supabase 암호화 저장",
            "  • GDPR 호환",
            "",
            "✅ 금융감독 준비",
            "  • Gemini AI 법률 자문 적용완료"
        ]
    )

    # 13. 연락처
    add_title_slide(
        "📞 연락처",
        "회사명: 수인AI브릿지\n\n"
        "대표자: 이희천\n"
        "연락처: 010-5650-0670\n"
        "이메일: hejunl@hanmail.net\n\n"
        "웹사이트: https://petcare-plus.vercel.app\n"
        "GitHub: https://github.com/HeeJeonLee/petcare-plus-v2\n\n"
        "🙏 감사합니다!"
    )

    # PPT 저장
    prs.save('/home/user/petcare-plus-v2/PetCare+_사업계획서.pptx')
    print("✅ PPT 생성 완료: /home/user/petcare-plus-v2/PetCare+_사업계획서.pptx")

if __name__ == '__main__':
    create_business_plan_ppt()
