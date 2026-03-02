#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
🎯 정책자금 신청용 최고의 사업계획서 생성기
- 모든 최종 수정안 반영
- 정책자금(정책금융) 신청에 최적화
- 전문성 + 완성도 극대화
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_premium_business_plan():
    """최고의 사업계획서 생성"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    BLUE = RGBColor(25, 118, 210)
    DARK_BLUE = RGBColor(13, 71, 161)
    LIGHT_BLUE = RGBColor(227, 242, 253)
    GREEN = RGBColor(56, 142, 60)
    ORANGE = RGBColor(251, 140, 0)
    GRAY = RGBColor(66, 66, 66)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    def add_title_slide(title, subtitle, company_name="수인AI브릿지"):
        """제목 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # 타이틀
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 부제
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(2)

        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        for text_line in subtitle.split('\n'):
            if subtitle_frame.paragraphs[0].text == "":
                p = subtitle_frame.paragraphs[0]
            else:
                p = subtitle_frame.add_paragraph()
            p.text = text_line
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(200, 220, 255)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(6)

        # 회사명 + 날짜
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9)
        height = Inches(0.5)

        company_box = slide.shapes.add_textbox(left, top, width, height)
        company_frame = company_box.text_frame
        p = company_frame.paragraphs[0]
        p.text = f"{company_name} | 2026년 3월"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 200, 255)
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title, contents_list, bg_color=LIGHT_GRAY):
        """내용 슬라이드 (제목 + 글머리)"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # 제목 바
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = BLUE
        title_shape.line.color.rgb = BLUE

        # 제목 텍스트
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = f"  📊 {title}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

        # 내용
        left = Inches(0.7)
        top = Inches(1.2)
        width = Inches(8.6)
        height = Inches(6)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for idx, item in enumerate(contents_list):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
            p.level = 0
            p.space_after = Pt(12)
            p.space_before = Pt(6)

    def add_two_column_slide(title, left_title, left_items, right_title, right_items):
        """2열 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_GRAY

        # 제목 바
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = BLUE
        title_shape.line.color.rgb = BLUE

        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = f"  📊 {title}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 왼쪽 열
        left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(4.5), Inches(6))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        p = left_frame.paragraphs[0]
        p.text = f"🎯 {left_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_after = Pt(10)

        for item in left_items:
            p = left_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
            p.level = 0

        # 오른쪽 열
        right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.2), Inches(4.5), Inches(6))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        p = right_frame.paragraphs[0]
        p.text = f"🎯 {right_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_after = Pt(10)

        for item in right_items:
            p = right_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.space_after = Pt(8)
            p.level = 0

    # ==================== 슬라이드 생성 시작 ====================

    # 1️⃣ 표지
    add_title_slide(
        "🐾 PetCare+",
        "AI 기반 펫보험 비교 플랫폼\n정책자금 신청 사업계획서"
    )

    # 2️⃣ 회사 개요
    add_content_slide(
        "회사 개요",
        [
            "📌 회사명: 수인AI브릿지",
            "📅 설립일: 2026년 2월 25일",
            "👤 대표: 이희천",
            "📍 위치: 경기도 수원시 영통구",
            "💼 사업자등록번호: 151-09-03201",
            "📞 연락처: 010-5650-0670 | hejunl@hanmail.net",
            "🌐 웹사이트: https://petcare-plus.vercel.app"
        ]
    )

    # 3️⃣ 비전과 미션
    add_two_column_slide(
        "비전 & 미션",
        "비전 (Vision)",
        [
            "🎯 대한민국 최고의 AI 기반",
            "    펫보험 비교 플랫폼",
            "",
            "💡 Claude AI를 활용한",
            "    혁신적 반려동물 보험 서비스",
            "",
            "🚀 차별화된 기술력으로",
            "    시장 주도"
        ],
        "미션 (Mission)",
        [
            "✅ 반려동물 보호자들에게",
            "    최적의 보험을 추천",
            "",
            "✅ 24시간 무료 AI 상담으로",
            "    접근성 극대화",
            "",
            "✅ 투명한 정보로",
            "    신뢰성 확보"
        ]
    )

    # 4️⃣ 핵심 기능 (1/2)
    add_content_slide(
        "핵심 기능 - Part 1️⃣",
        [
            "🤖 AI 맞춤형 펫보험 추천",
            "   • Claude API 기반 개인화 추천",
            "   • 펫의 나이, 종, 질병력 분석",
            "",
            "📊 8개 보험사 실시간 비교",
            "   • 메리츠, 삼성, DB, 현대, KB, 한화, 농협, 롯데",
            "   • 보험료, 보장범위, 특약 상세 비교",
            "",
            "💬 24시간 AI 챗봇 상담",
            "   • Claude API 기반 자연스러운 대화",
            "   • 실시간 질문 답변 (언제 어디서나)"
        ]
    )

    # 5️⃣ 핵심 기능 (2/2)
    add_content_slide(
        "핵심 기능 - Part 2️⃣",
        [
            "🏥 주변 동물병원 검색",
            "   • Google Maps API 기반 위치 검색",
            "   • 가까운 병원 정보 실시간 제공",
            "   • 병원 리뷰 및 평점 표시",
            "",
            "📋 보험금 청구 프로세스 가이드",
            "   • 단계별 청구 절차 설명",
            "   • 필요 서류 체크리스트",
            "",
            "📧 AI 맞춤형 상담 신청",
            "   • 고객정보 수집 및 이메일 발송 (Resend API)",
            "   • 전문가 연결로 전환율 극대화"
        ]
    )

    # 6️⃣ 기술 스택
    add_two_column_slide(
        "기술 스택",
        "프론트엔드",
        [
            "⚛️ React 18 (UI 라이브러리)",
            "⚡ Vite (빌드 도구)",
            "🎨 Tailwind CSS (스타일링)",
            "📱 PWA (모바일 앱처럼 사용)",
            "📊 React Hooks (상태 관리)",
            "",
            "성능: 초로딩 < 2초, 모바일 최적화"
        ],
        "백엔드 & 클라우드",
        [
            "🤖 Claude API (AI 챗봇)",
            "🗄️ Supabase (데이터베이스)",
            "📧 Resend API (이메일 발송)",
            "🗺️ Google Maps API (병원검색)",
            "🚀 Vercel (호스팅)",
            "⚙️ Node.js (API 서버)",
            "",
            "장점: 서버리스, 비용 효율적"
        ]
    )

    # 7️⃣ 시장 분석
    add_two_column_slide(
        "시장 분석",
        "시장 기회",
        [
            "📈 국내 펫보험 시장 규모",
            "   2024: 약 3,000억 원",
            "   2025: 약 4,500억 원 (예상)",
            "",
            "🐕 반려동물 보유 가구",
            "   약 1,000만 가구 (전국 기준)",
            "",
            "💰 평균 월 보험료",
            "   25,000~35,000원 (정상 수준)",
            "",
            "✅ 성장률: 연 30~50%"
        ],
        "경쟁 우위",
        [
            "🤖 AI 기반 자동화",
            "   • 인간 상담사 불필요",
            "   • 24시간 자동 운영",
            "",
            "💎 차별성",
            "   • 유일한 AI 맞춤형 추천",
            "   • 정책금융 지원으로 신뢰성",
            "",
            "🚀 확장성",
            "   • 자율진화형 플랫폼",
            "   • 자동 콘텐츠 생성",
            "   • SNS 자동 마케팅"
        ]
    )

    # 8️⃣ 수익 모델
    add_content_slide(
        "수익 모델 & 예상 매출",
        [
            "💰 Phase 1 (2026.03-05): 기반 수립",
            "   • 제휴 수수료: 보험료의 3~5% (보험사별)",
            "   • 예상 월 매출: 500만원 (가입자 200명 기준)",
            "",
            "💰 Phase 2 (2026.06-12): 성장",
            "   • 광고 수익: 동물병원, 펫용품업체",
            "   • 프리미엄 상담 (선택사항)",
            "   • 예상 월 매출: 2,000~3,000만원",
            "",
            "💰 Phase 3 (2027년~): 고도화",
            "   • 보험료 기반 수익: 5~10억원/년 (목표)",
            "   • SNS 광고, 콘텐츠 마케팅",
            "   • B2B 기업 제휴 (펫케어 업체)"
        ]
    )

    # 9️⃣ 자율진화형 플랫폼
    add_content_slide(
        "혁신: 자율진화형 플랫폼",
        [
            "🔄 핵심: Claude Code를 활용한 자동화",
            "   • 90% 이상 자동화 (인간 개입 최소화)",
            "",
            "📊 데이터 자동 갱신",
            "   • 보험료 변동 자동 감지",
            "   • 새 상품 자동 추가",
            "   • 고객에게 자동 알림",
            "",
            "✍️ 콘텐츠 자동 생성",
            "   • 일일 블로그 포스트 (Claude API)",
            "   • SNS 콘텐츠 자동 작성 (Instagram/Twitter/Facebook)",
            "   • 주간 뉴스레터 자동화",
            "",
            "🎯 AI 기반 최적화",
            "   • 사용자 행동 분석",
            "   • UI/UX 자동 개선",
            "   • A/B 테스트 자동 실행"
        ]
    )

    # 🔟 팀 구성
    add_two_column_slide(
        "팀 구성",
        "현재 (2026년 3월)",
        [
            "👤 이희천",
            "   • CEO / 사업가",
            "   • 보험 컨설턴트",
            "   • 전략 수립",
            "",
            "🤖 Claude Code",
            "   • 개발 자동화",
            "   • AI 기반 의사결정",
            "   • 콘텐츠 생성",
            "",
            "구성: 1명 + AI"
        ],
        "향후 확대안 (2026년 4분기)",
        [
            "💼 영업 담당 (1명)",
            "   • 보험사 제휴",
            "   • B2B 계약",
            "",
            "👨‍💻 개발자 (2-3명)",
            "   • 모바일 앱 개발",
            "   • 백엔드 최적화",
            "",
            "🎨 UX/UI 디자이너 (1명)",
            "   • 사용자 경험 개선",
            "",
            "구성: 5-6명 + AI (예상)"
        ]
    )

    # 1️⃣1️⃣ 개발 로드맵
    add_content_slide(
        "개발 로드맵 (2026)",
        [
            "🟢 Phase 1: 초기 구축 (2026.02.25 ~ 3월)",
            "   ✅ 완료: 기본 기능 (AI 추천, 비교표, 챗봇, 병원검색)",
            "   ✅ 완료: Vercel 배포 및 API 키 환경변수 설정",
            "   ✅ 완료: 법적 준칙 적용 (면책공고)",
            "",
            "🟡 Phase 2: 자율진화 구축 (2026.04 ~ 5월)",
            "   📊 데이터 자동 갱신 시스템",
            "   ✍️ 콘텐츠 자동 생성 시스템",
            "   📈 Analytics + AI 최적화 모듈",
            "   💻 관리자 대시보드 개발",
            "",
            "🟠 Phase 3: 고도화 (2026.06 ~ 12월)",
            "   📱 iOS/Android 네이티브 앱",
            "   🌍 국제 확장 (영문 지원)",
            "   💳 결제 시스템 통합",
            "   🎯 AI 기반 마케팅 자동화"
        ]
    )

    # 1️⃣2️⃣ 법적 준칙 & 규제
    add_content_slide(
        "법적 준칙 & 규제 준수",
        [
            "⚖️ 보험업법 준수",
            "   • 보험 상품 권유 아님 (정보 제공만 명시)",
            "   • 면책공고 3곳 배치 (중요 고지 사항)",
            "   • 확정적 표현 제거 (\"최저가\", \"반드시\")",
            "",
            "🔐 개인정보보호법 (PIPA)",
            "   • 고객 정보 암호화 저장 (Supabase)",
            "   • GDPR 준비 (국제 확장 대비)",
            "   • 개인정보 처리방침 명시",
            "",
            "💰 금융감독 (FSS)",
            "   • 정책금융 신청 시 별도 심의",
            "   • 금융감독당국 보고 체계",
            "",
            "📝 사전 심의 제외",
            "   • 정보 제공 형식으로 법적 안전성 확보"
        ]
    )

    # 1️⃣3️⃣ 정책자금 신청 전략
    add_content_slide(
        "정책자금 신청 전략",
        [
            "🎯 신청 대상",
            "   • 정책금융(중소기업진흥공단 등)",
            "   • 스타트업 지원금",
            "   • 혁신성장 자금",
            "",
            "💡 차별성 강조",
            "   ✅ AI 기반 혁신 서비스 (Claude API)",
            "   ✅ 자율진화형 플랫폼 (자동화 90%)",
            "   ✅ 고성장 시장 (연 30~50% 성장)",
            "   ✅ 낮은 초기 자본금 (Vercel 활용)",
            "",
            "📊 사용처",
            "   • 개발 인력 충원 (2-3명): 약 1억원",
            "   • 마케팅 비용: 약 3,000만원",
            "   • 인프라 및 운영비: 약 2,000만원",
            "   • 총 신청액: 약 1.5억원 (예상)"
        ]
    )

    # 1️⃣4️⃣ 마지막 슬라이드 (연락처)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "감사합니다! 🙏"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 연락처 정보
    info_box = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(3))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True

    contact_info = [
        ("회사명", "수인AI브릿지"),
        ("대표자", "이희천"),
        ("이메일", "hejunl@hanmail.net"),
        ("전화", "010-5650-0670"),
        ("웹사이트", "https://petcare-plus.vercel.app"),
        ("GitHub", "https://github.com/HeeJeonLee/petcare-plus-v2"),
    ]

    for idx, (label, value) in enumerate(contact_info):
        if idx > 0:
            info_frame.add_paragraph()
        p = info_frame.paragraphs[idx]
        p.text = f"📌 {label}: {value}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.space_before = Pt(10)
        p.space_after = Pt(10)

    # PPT 저장
    output_path = "/home/user/petcare-plus-v2/PetCare+_정책자금신청용_사업계획서.pptx"
    prs.save(output_path)

    print(f"✅ 사업계획서 생성 완료!")
    print(f"📄 저장 위치: {output_path}")
    print(f"📊 슬라이드 수: {len(prs.slides)}")

    # 파일 크기 확인
    file_size = os.path.getsize(output_path)
    print(f"📦 파일 크기: {file_size / 1024:.1f} KB")

    return output_path

if __name__ == "__main__":
    create_premium_business_plan()
