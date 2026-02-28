#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PetCare+ 사업계획서 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    BLUE = RGBColor(37, 99, 235)
    PURPLE = RGBColor(124, 58, 237)
    DARK_GRAY = RGBColor(30, 41, 59)
    LIGHT_GRAY = RGBColor(248, 250, 252)
    WHITE = RGBColor(255, 255, 255)

    def add_title_slide(title, subtitle):
        """제목 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLUE

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 부제
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_list):
        """내용 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # 제목
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY

        # 선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0))
        line.line.color.rgb = BLUE
        line.line.width = Pt(3)

        # 내용
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_list):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(8)
            p.space_after = Pt(8)

        return slide

    # 1. 표지
    add_title_slide("🐾 PetCare+", "AI 기반 펫보험 비교 플랫폼")

    # 2. 기본 정보
    add_content_slide("📋 기본 정보", [
        "회사명: 수인AI브릿지",
        "사업자등록번호: 151-09-03201",
        "개업일: 2026년 02월 25일",
        "대표자: 이희전",
        "사업분야: AI 기반 펫보험 비교 및 상담 플랫폼",
        "신청자금: 5천만원"
    ])

    # 3. 사업 개요
    add_content_slide("1️⃣ 사업 개요", [
        "Claude AI 기반의 자동화된 펫보험 비교 및 24시간 상담 플랫폼",
        "",
        "📊 시장 현황:",
        "• 반려동물 양육 가구: 600만 가구 (2024)",
        "• 펫보험 시장: 5,000억원 (연평균 15-20% 성장)",
        "• 펫보험 가입률: 5% (미국 20%에 비해 저조)",
    ])

    # 4. 문제점과 솔루션
    add_content_slide("🎯 문제점 vs 솔루션", [
        "❌ 문제점:",
        "• 8개 이상 보험사의 복잡한 상품 비교 어려움",
        "• 고가의 전문 상담 비용 (시간당 5~10만원)",
        "• 야간/휴일 상담 불가능",
        "",
        "✅ 우리의 솔루션:",
        "• AI 기반 자동 비교 시스템",
        "• 24시간 무료 상담 (Claude AI)",
        "• 완전 자동화된 운영"
    ])

    # 5. 기술 스택 - AI 챗봇
    add_content_slide("2️⃣ 핵심 기술 - AI 챗봇", [
        "기술 스택: Claude API (Sonnet 4.6, Opus 4.6)",
        "",
        "✨ 구현 내용:",
        "• Claude API 통합: 실시간 자연어 처리",
        "• System Prompt 엔지니어링: 펫보험 전문가 페르소나",
        "• Multi-turn 대화: 견종별/나이별 맞춤 분석",
        "",
        "💰 효과:",
        "• 상담사 인건비 제거 (월 500만원 절감)",
        "• 1000명 동시 상담 가능"
    ])

    # 6. 기술 스택 - 분석 & 최적화
    add_content_slide("🔄 자동 성능 분석 & 최적화", [
        "기술 스택: Custom Analytics Module, Claude AI",
        "",
        "📊 실시간 분석:",
        "• 클릭 패턴, 스크롤 깊이, 세션 지속시간 추적",
        "• AI 기반 자동 최적화 제안",
        "• A/B 테스트 자동 실행",
        "",
        "목표: 전환율 30% 증가"
    ])

    # 7. 콘텐츠 자동화
    add_content_slide("📝 AI 콘텐츠 자동 생성", [
        "기술 스택: Claude API, Scheduling System",
        "",
        "자동화 대상:",
        "• 일일 블로그 포스트 (1500-2000자)",
        "• SNS 콘텐츠 (Instagram, Twitter, Facebook)",
        "",
        "효과:",
        "• 콘텐츠 마케팅 비용 제거 (월 200만원 절감)",
        "• 월 50개 이상 콘텐츠 생성",
        "• 유기 트래픽 35% 증가 예상"
    ])

    # 8. 시장 분석
    add_content_slide("3️⃣ 시장 분석", [
        "📈 시장 규모:",
        "• 펫보험 시장: 5,000억원 (2024)",
        "• 연평균 성장률: 15-20%",
        "• 2028년 예상: 약 10,000억원",
        "",
        "🎯 타겟 고객:",
        "• 반려동물 보유 가구: 600만 가구",
        "• 정보 수집 관심층: 150만 가구 (25%)",
        "• 우리의 TAM: 150만 명 (1차), 50만 명 (3년차)"
    ])

    # 9. 경쟁 분석
    add_content_slide("⚔️ 경쟁사 비교", [
        "기존 플레이어:",
        "• 보험사 홈페이지: 복잡, 상담 유료",
        "• 비교 사이트: 수동 상담, AI 없음",
        "• 직거래 상담: 비용 높음, 시간 제한",
        "",
        "🌟 PetCare+ 차별점:",
        "✅ Claude AI 기반 24시간 무료 상담",
        "✅ 자동화로 초저가 모델",
        "✅ 자동 콘텐츠 마케팅"
    ])

    # 10. 재무 계획
    add_content_slide("4️⃣ 재무 계획", [
        "💰 초기 투자 5,000만원:",
        "• 개발 비용: 1,000만원",
        "• 마케팅 비용: 1,500만원",
        "• 인프라/운영: 1,500만원",
        "• 인건비: 1,000만원",
        "",
        "월 고정 운영비: 100만원",
        "(기존 모델 3,000만원 대비 90% 절감)"
    ])

    # 11. 예상 수익
    add_content_slide("📊 3년 매출 전망", [
        "Year 1 (2026): 1,500만원 (투자 대비 3배)",
        "• 월 방문자: 2,000명 → 10,000명",
        "• 월 계약: 100건 → 200건",
        "",
        "Year 2 (2027): 1억 2,000만원",
        "• 월 계약: 200건",
        "",
        "Year 3 (2028): 3억원",
        "• 월 계약: 500건",
        "",
        "손익분기점: 6개월"
    ])

    # 12. 실행 계획
    add_content_slide("5️⃣ 실행 계획", [
        "Phase 1 (1월-3월): 기반 조성 ✅ 완료",
        "• AI 챗봇, 자동화 시스템 개발",
        "",
        "Phase 2 (4월-6월): 마켓 진출",
        "• 목표: 월 2,000명, 100건 상담신청",
        "• 마케팅 비용: 1,300만원",
        "",
        "Phase 3 (7월-12월): 성장 단계",
        "• 목표: 월 10,000명, 500건 상담",
        "• ROI 150% 달성"
    ])

    # 13. 성공 지표
    add_content_slide("🎯 성공 지표 (KPI)", [
        "📍 6개월 목표:",
        "• 월간 방문자: 2,000명 이상",
        "• 상담 신청률: 5% 이상",
        "• 고객 만족도: 4.5/5.0 이상",
        "",
        "📍 1년 목표:",
        "• 월간 방문자: 10,000명 이상",
        "• 월간 계약: 200건 이상",
        "• 손익분기점 달성",
        "",
        "📍 3년 목표:",
        "• 월간 방문자: 50,000명 이상",
        "• 연간 수익: 3억원 이상"
    ])

    # 14. 핵심 가치
    add_content_slide("💎 핵심 가치", [
        "✅ 기술 혁신: AI 기반 자동화로 90% 비용 절감",
        "",
        "✅ 시장 타이밍: 펫보험 시장 급성장 초입",
        "",
        "✅ 비즈니스 모델: 검증된 수수료 기반 수익구조",
        "",
        "✅ 확장성: 자동화로 무한 스케일 가능",
        "",
        "✅ 지속성: 자율진화로 경쟁력 자동 유지"
    ])

    # 15. 정책자금 신청 사유
    add_content_slide("🏛️ 정책자금 신청 사유", [
        "1️⃣ 기술 혁신: Claude AI 활용한 K-스타트업 모델",
        "",
        "2️⃣ 일자리 창출: 자동화로 양질의 일자리 창출",
        "",
        "3️⃣ 시장 선도: 펫보험 시장에서 AI 선도 업체",
        "",
        "4️⃣ 수출 가능성: 자동화 시스템으로 글로벌 확장 용이",
        "",
        "5️⃣ 경제 파급: 보험 중개 시장 활성화"
    ])

    # 16. 마지막 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PURPLE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "감사합니다! 🐾"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "수인AI브릿지 | 2026년 02월 28일"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # PPT 저장
    output_path = "/home/user/petcare-plus-v2/PetCare+_사업계획서.pptx"
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
